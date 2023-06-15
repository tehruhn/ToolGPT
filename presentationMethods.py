import collections 
import collections.abc

import re
import openai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.dml.color import RGBColor

COMPLETION_MODEL = "gpt-3.5-turbo-0613"

def create_presentation(filename):
    """
    Creates a new PowerPoint presentation

    Parameters
    ----------
    filename : string
        The name of the PowerPoint file to be created

    Returns
    -------
    idk
    """
    prs = Presentation()
    prs.save(filename)
    print(f"Presentation '{filename}' created.")

def add_slide_with_bullets(filename, title, bullet_points_text):
    """
    Adds a slide with bullet points to an existing PowerPoint presentation

    Parameters
    ----------
    filename : string
        The name of the PowerPoint file to add a slide to
    title : string
        The title of the slide
    bullet_points_text : string
        The bullet points for the slide, separated by newline

    Returns
    -------
    idk
    """
    bullet_points = bullet_points_text.split("\\n")
    bullet_points = [elem.strip() for elem in bullet_points]
    prs = Presentation(filename)

    slide_layout = prs.slide_layouts[5]  # blank slide
    slide = prs.slides.add_slide(slide_layout)

    # define title box
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    tf.text_anchor = MSO_ANCHOR.MIDDLE

    p = tf.add_paragraph()
    p.text = title
    p.font.bold = True
    p.font.size = Pt(36)
    p.font.color.rgb = RGBColor(0x00, 0x00, 0x80)  # Dark blue color
    p.alignment = PP_ALIGN.CENTER

    # define content box with bullet points
    left = Inches(1)
    top = Inches(2)
    width = Inches(8.5)
    height = Inches(5)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for point in bullet_points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0  # making the text a bullet point
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(0x80, 0x80, 0x80)  # Grey color

    prs.save(filename)
    print(f"Slide added to the presentation '{filename}'.")


def get_completion(messages):
    response = openai.ChatCompletion.create(
        model=COMPLETION_MODEL,
        messages=messages,
        functions=[
            {
                "name": "create_presentation",
                "description": "Creates a new PowerPoint presentation",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "filename": {
                            "type": "string",
                            "description": "The name of the PowerPoint file to be created",
                        },
                    },
                    "required": ["filename"],
                },
            },
            {
                "name": "add_slide_with_bullets",
                "description": "Adds a slide with bullet points to an existing PowerPoint presentation",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "filename": {
                            "type": "string",
                            "description": "The name of the PowerPoint file to add a slide to",
                        },
                        "title": {
                            "type": "string",
                            "description": "The title of the slide",
                        },
                        "bullet_points_text": {
                            "type": "string",
                            "description": "The bullet points for the slide, separated by periods",
                        },
                    },
                    "required": ["filename", "title", "bullet_points_text"],
                },
            },
        ],
        temperature=0,
    )

    return response
